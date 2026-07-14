from __future__ import annotations

import json
from pathlib import Path
from typing import List

from tangyuan.tools.paths import workspace_path

def create_pptx(workspace: Path, rel: str, title: str, slides: List[str]) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return json.dumps({"ok": False, "error": "缺少 python-pptx，请 pip install python-pptx"}, ensure_ascii=False)

    path = workspace_path(workspace, rel)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "由汤圆 Tangyuan 生成"

    for raw in slides:
        lines = [x.strip() for x in str(raw).splitlines() if x.strip()]
        if not lines:
            continue
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = lines[0][:120]
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        first = True
        for line in lines[1:] or [" "]:
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = line
            p.level = 0
            p.font.size = Pt(18)

    prs.save(str(path))
    return json.dumps({"ok": True, "path": str(path), "slides": len(slides) + 1}, ensure_ascii=False)

